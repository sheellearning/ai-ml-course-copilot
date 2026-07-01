from typing import List, Dict, Any
from pptx import Presentation
import os


# Load text from each PowerPoint slide with source metadata.
def load_pptx_file(pptx_path: str) -> List[Dict[str, Any]]:
    presentation = Presentation(pptx_path)
    records = []
    filename = os.path.basename(pptx_path)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []

        # Extract readable text from all slide shapes.
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text_parts.append(text)

        slide_text = "\n".join(slide_text_parts).strip()

        if not slide_text:
            continue

        metadata = {
            "source": "course",
            "doc_type": "pptx",
            "filename": filename,
            "slide": slide_index,
            "path": pptx_path,
        }

        records.append(
            {
                "text": slide_text,
                "metadata": metadata,
            }
        )

    return records