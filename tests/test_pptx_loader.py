from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation

from core.loaders.pptx_loader import load_pptx_file


# Create a temporary PowerPoint presentation for testing.
def create_sample_pptx(pptx_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    slide.shapes.title.text = "Large Language Models"
    slide.placeholders[1].text = "LLMs can process and generate natural language."

    presentation.save(pptx_path)


# Verify that slide text and metadata are extracted correctly.
def test_pptx_loader():
    with TemporaryDirectory() as temporary_folder:
        pptx_path = Path(temporary_folder) / "sample.pptx"
        create_sample_pptx(pptx_path)

        records = load_pptx_file(str(pptx_path))

        assert len(records) == 1

        record = records[0]

        assert "Large Language Models" in record["text"]
        assert "LLMs can process and generate natural language." in record["text"]
        assert record["metadata"]["source"] == "course"
        assert record["metadata"]["doc_type"] == "pptx"
        assert record["metadata"]["filename"] == "sample.pptx"
        assert record["metadata"]["slide"] == 1
        assert record["metadata"]["path"] == str(pptx_path)


if __name__ == "__main__":
    test_pptx_loader()
    print("PPTX LOADER TEST PASSED")
